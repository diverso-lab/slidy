from pptx import Presentation
import markdown
from bs4 import BeautifulSoup
from copy import deepcopy
from PIL import Image


class PPTXGenerator:
    INDEX_SLIDE_POSITION = 1
    CONTENT_SLIDE_POSITION = 2

    def __init__(self, template_path):
        self.prs = Presentation(template_path)

    def duplicate_slide(self, slide_index):
        """Duplica una diapositiva en una presentación."""
        source = self.prs.slides[slide_index]
        slide_layout = source.slide_layout
        new_slide = self.prs.slides.add_slide(slide_layout)

        for shape in reversed(source.shapes):
            el = shape.element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert(2, new_el)

        return new_slide

    @staticmethod
    def replace_text_while_keeping_format(shape, placeholder, new_text):
        """Reemplaza texto en una forma mientras mantiene el formato original."""
        if placeholder in shape.text:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, new_text)

    @staticmethod
    def add_image_within_bounds(slide, image_path, sample_shape):
        """Agrega una imagen a una diapositiva sin exceder las dimensiones de una forma de muestra."""
        # Obtener las dimensiones de la forma de muestra
        max_width = sample_shape.width
        max_height = sample_shape.height

        # Obtener las dimensiones de la imagen
        with Image.open(image_path) as img:
            img_width, img_height = img.size

        # Calcular la relación de aspecto de la imagen
        aspect_ratio = img_width / img_height

        # Determinar las dimensiones de la imagen ajustada
        if max_width / max_height > aspect_ratio:
            new_height = max_height
            new_width = max_height * aspect_ratio
        else:
            new_width = max_width
            new_height = max_width / aspect_ratio

        # Añadir la imagen a la diapositiva con las dimensiones ajustadas
        left = sample_shape.left
        top = sample_shape.top
        pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)

        return pic

    @staticmethod
    def get_sample_shape(slide, sample_name="SAMPLE_IMAGE"):
        """Obtiene la forma de muestra por su nombre."""
        for shape in slide.shapes:
            if shape.name == sample_name:
                return shape
        return None

    def duplicate_and_highlight_index(self, sections, current_section):
        """Duplica el slide del índice, actualiza los títulos de las secciones y resalta la sección actual."""
        slide = self.duplicate_slide(self.INDEX_SLIDE_POSITION)

        # Buscamos el shape que contiene el identificador #item_list
        for shape in slide.shapes:
            if shape.has_text_frame:
                paragraphs = shape.text_frame.paragraphs
                for paragraph in paragraphs:
                    if "#item_list" in paragraph.text:
                        # Guardamos las propiedades originales del run
                        original_run = paragraph.runs[0]
                        original_font_name = original_run.font.name
                        original_font_size = original_run.font.size
                        try:
                            original_font_color = original_run.font.color.rgb
                        except AttributeError:
                            original_font_color = None  # Si el color no se define en términos de RGB

                        # Eliminamos el placeholder original
                        shape.text_frame._element.remove(paragraph._element)

                        # Agregamos cada sección con el estilo original
                        for index, section in enumerate(sections, 1):  # index empieza en 1
                            new_paragraph = shape.text_frame.add_paragraph()
                            new_run = new_paragraph.add_run()
                            new_run.text = f"{index}. {section}"
                            new_run.font.bold = True if section == current_section else False
                            new_run.font.name = original_font_name
                            new_run.font.size = original_font_size
                            if original_font_color:
                                new_run.font.color.rgb = original_font_color

        return slide

    def process_header(self, soup):
        title = soup.find('h1').get_text()
        authors_line = soup.find('h2').get_text()
        workshop = soup.find('h3').get_text()

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self.replace_text_while_keeping_format(shape, "#main_title", title)
                    self.replace_text_while_keeping_format(shape, "#authors",
                                                           authors_line.replace("Authors:", "").strip())
                    self.replace_text_while_keeping_format(shape, "#workshop", workshop)

    def process_subsection_text(self, slide, subtitle, section, section_counter, subtitle_counter,
                                total_content_slides, content_slide_counter):
        for shape in slide.shapes:
            if shape.has_text_frame:
                section_title = f"{section_counter}. {section.get_text()}"
                section_subtitle = f"{section_counter}.{subtitle_counter}. {subtitle.get_text()}"
                slide_number_format = f"{content_slide_counter:02} - {total_content_slides:02}"

                content_text = subtitle.find_next('p').get_text() if subtitle.find_next('p') else ""
                self.replace_text_while_keeping_format(shape, "#content", content_text)
                self.replace_text_while_keeping_format(shape, "#section_title", section_title)
                self.replace_text_while_keeping_format(shape, "#section_subtitle", section_subtitle)
                self.replace_text_while_keeping_format(shape, "#slide_number", slide_number_format)

    def process_subsection_image(self, slide, subtitle):
        sample_shape = self.get_sample_shape(slide)
        image_tag = subtitle.find_next('img')
        if sample_shape:
            if image_tag:
                image_path = image_tag['src']
                self.add_image_within_bounds(slide, image_path, sample_shape)
            sp = sample_shape._element
            slide.shapes._spTree.remove(sp)

    def process_section(self, section, sections, section_counter, total_content_slides, content_slide_counter):
        self.duplicate_and_highlight_index([s.get_text() for s in sections], section.get_text())
        following_elements = section.find_all_next()
        subtitles = []
        for elem in following_elements:
            if elem.name == 'h3':
                subtitles.append(elem)
            elif elem.name == 'h2':
                break

        subtitle_counter = 1
        for subtitle in subtitles:
            slide = self.duplicate_slide(self.CONTENT_SLIDE_POSITION)
            self.process_subsection_text(slide, subtitle, section, section_counter, subtitle_counter,
                                         total_content_slides, content_slide_counter)
            self.process_subsection_image(slide, subtitle)
            subtitle_counter += 1
            content_slide_counter += 1
        return content_slide_counter

    def generate_from_markdown(self, input_file):
        with open(input_file, 'r') as f:
            content = f.read()
        html_content = markdown.markdown(content)
        soup = BeautifulSoup(html_content, 'html.parser')

        self.process_header(soup)

        total_content_slides = len(soup.find_all('h3')) - 1
        sections = soup.find_all('h2')[1:]

        section_counter = 1
        content_slide_counter = 1
        for section in sections:
            content_slide_counter = self.process_section(section, sections, section_counter, total_content_slides,
                                                         content_slide_counter)

            section_counter += 1

        slide_id = self.prs.slides._sldIdLst[1].rId
        self.prs.part.drop_rel(slide_id)
        del self.prs.slides._sldIdLst[1]
        del self.prs.slides._sldIdLst[1]

        output_path = "output.pptx"
        self.prs.save(output_path)
        print(f"Presentación guardada en {output_path}")
