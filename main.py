from pptx import Presentation
import markdown
from bs4 import BeautifulSoup
from copy import deepcopy


def duplicate_slide(prs, slide_index):
    """Duplica una diapositiva en una presentación."""
    source = prs.slides[slide_index]
    slide_layout = source.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    for shape in reversed(source.shapes):
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert(2, new_el)

    return new_slide


def replace_text_while_keeping_format(shape, placeholder, new_text):
    """Reemplaza texto en una forma mientras mantiene el formato original."""
    if placeholder in shape.text:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, new_text)


def duplicate_and_highlight_index(prs, sections, current_section):
    """Duplica el slide del índice, actualiza los títulos de las secciones y resalta la sección actual."""
    slide = duplicate_slide(prs, 1)  # Duplicamos el tercer slide (índice)

    # Buscamos el shape que contiene el identificador #item_list
    for shape in slide.shapes:
        if shape.has_text_frame:
            paragraphs = shape.text_frame.paragraphs
            for paragraph in paragraphs:
                if "#item_list" in paragraph.text:
                    # Guardamos las propiedades originales del run
                    original_run = paragraph.runs[0]
                    original_font_name = original_run.font.name
                    original_font_size = original_run.font.size
                    try:
                        original_font_color = original_run.font.color.rgb
                    except AttributeError:
                        original_font_color = None  # Si el color no se define en términos de RGB

                    # Eliminamos el placeholder original
                    shape.text_frame._element.remove(paragraph._element)

                    # Agregamos cada sección con el estilo original
                    for index, section in enumerate(sections, 1):  # index empieza en 1
                        new_paragraph = shape.text_frame.add_paragraph()
                        new_run = new_paragraph.add_run()
                        new_run.text = f"{index}. {section}"
                        new_run.font.bold = True if section == current_section else False
                        new_run.font.name = original_font_name
                        new_run.font.size = original_font_size
                        if original_font_color:
                            new_run.font.color.rgb = original_font_color

    return slide


def markdown_to_pptx_with_template(input_file, template_path):
    with open(input_file, 'r') as f:
        content = f.read()
    html_content = markdown.markdown(content)
    soup = BeautifulSoup(html_content, 'html.parser')

    title = soup.find('h1').get_text()
    authors_line = soup.find('h2').get_text()
    workshop = soup.find('h3').get_text()

    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                replace_text_while_keeping_format(shape, "#main_title", title)
                replace_text_while_keeping_format(shape, "#authors", authors_line.replace("Authors:", "").strip())
                replace_text_while_keeping_format(shape, "#workshop", workshop)

    total_content_slides = sum(1 for section in soup.find_all('h2')[1:] for _ in section.find_all_next('h3'))

    sections = soup.find_all('h2')[1:]
    section_titles = [section.get_text() for section in sections]
    section_counter = 1
    content_slide_counter = 0

    for section in sections:
        # Creamos el slide del índice y lo resaltamos para la sección actual
        duplicate_and_highlight_index(prs, [s.get_text() for s in sections], section.get_text())
        # (aquí va el resto del código para crear las diapositivas de contenido)

        subtitles = section.find_all_next('h3')
        subtitle_counter = 1
        for subtitle in subtitles:
            slide = duplicate_slide(prs, 2)
            content_slide_counter += 1
            for shape in slide.shapes:
                if shape.has_text_frame:
                    section_title = f"{section_counter}. {section.get_text()}"
                    section_subtitle = f"{section_counter}.{subtitle_counter}. {subtitle.get_text()}"
                    slide_number_format = f"{content_slide_counter:02} - {total_content_slides:02}"

                    replace_text_while_keeping_format(shape, "#section_title", section_title)
                    replace_text_while_keeping_format(shape, "#section_subtitle", section_subtitle)
                    replace_text_while_keeping_format(shape, "#main_title", title)
                    replace_text_while_keeping_format(shape, "#slide_number", slide_number_format)
            subtitle_counter += 1
        section_counter += 1

    slide_id = prs.slides._sldIdLst[1].rId
    prs.part.drop_rel(slide_id)
    del prs.slides._sldIdLst[1]
    del prs.slides._sldIdLst[1]


    output_path = "output.pptx"
    prs.save(output_path)
    print(f"Presentación guardada en {output_path}")


if __name__ == "__main__":
    markdown_file = "ejemplo_presentacion.md"
    template_file = "plantilla.pptx"
    markdown_to_pptx_with_template(markdown_file, template_file)
